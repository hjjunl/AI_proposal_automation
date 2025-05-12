# 재실행을 위한 필요 모듈 재로드 + 경로 반영 버전 통합
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
import matplotlib.pyplot as plt
import os

# 디렉터리 상수
TEMPLATE_DIR = "DB/ppt_template"
OUTPUT_DIR = "DB/proposal_result"
RFP_DIR = "DB/rfp"

SLIDE_LAYOUT_MAP = {
    "cover_page": "Title Slide",
    "table_of_contents": "Title and Content",
    "project_understanding": "Title and Content",
    "client_needs_summary": "Title and Content",
    "market_analysis_market_overview": "Title and Content",
    "growth_trend_analysis": "Title and Content",
    "industry_drivers_challenges": "Title and Content",
    "competitive_benchmarking": "Title and Content",
    "swot_analysis": "Title and Content",
    "solution_overview": "Title and Content",
    "strategic_recommendations": "Title and Content",
    "implementation_plan": "Title and Content",
    "timeline_milestones": "Title and Content",
    "risk_management_plan": "Title and Content",
    "expected_benefits": "Title and Content",
    "investment_budget_estimation": "Title and Content",
    "team_introduction": "Title and Content",
    "why_us_differentiation": "Title and Content",
    "closing_summary": "Title and Content",
    "qna": "Title Only"
}

def get_layout_by_name(prs, layout_name):
    for layout in prs.slide_layouts:
        if layout.name == layout_name:
            return layout
    return prs.slide_layouts[0]

def insert_graph(slide, data):
    fig, ax = plt.subplots()
    ax.plot(data['x'], data['y'], marker='o')
    ax.set_title(data.get('title', '그래프'))
    graph_path = "temp_graph.png"
    fig.savefig(graph_path)
    plt.close(fig)
    slide.shapes.add_picture(graph_path, Inches(1), Inches(2), Inches(6), Inches(3))
    os.remove(graph_path)

def generate_slide_content_with_llm(slide_key, context="양자컴퓨팅 시장"):
    if slide_key == "market_analysis_market_overview":
        return {
            "title": "시장 개요",
            "text": f"{context}은 향후 5년간 연평균 성장률 15% 이상으로 고성장할 것으로 예상됩니다.",
            "graph": {
                "x": [2021, 2022, 2023, 2024, 2025],
                "y": [1.2, 1.8, 2.4, 3.1, 3.9],
                "title": f"{context} 시장 규모 추이"
            }
        }
    elif slide_key == "competitive_benchmarking":
        return {
            "title": "경쟁사 비교 분석",
            "text": "주요 경쟁사 대비 당사의 기술력 및 전략적 포지셔닝을 정리합니다.",
        }
    return {"title": slide_key, "text": f"{slide_key}에 대한 설명입니다."}

def create_presentation_from_template(template_path):
    return Presentation(template_path)

def add_slide_from_template(prs, slide_key, slide_data, template_path):
    layout_name = SLIDE_LAYOUT_MAP.get(slide_key, "Title and Content")
    base = Presentation(template_path)
    layout = get_layout_by_name(base, layout_name)
    slide = prs.slides.add_slide(layout)

    if slide.shapes.title and 'title' in slide_data:
        slide.shapes.title.text = slide_data['title']

    if 'text' in slide_data:
        for shape in slide.shapes:
            if shape.is_placeholder and shape.placeholder_format.idx == 1:
                shape.text = slide_data['text']
                break

    if 'graph' in slide_data:
        insert_graph(slide, slide_data['graph'])

def generate_ppt_from_template(slide_keys, template_path, output_path):
    prs = create_presentation_from_template(template_path)
    for slide_key in slide_keys:
        slide_data = generate_slide_content_with_llm(slide_key)
        add_slide_from_template(prs, slide_key, slide_data, template_path)
    prs.save(output_path)
    return output_path

# 실제 경로를 이용한 실행
template_path = f"{TEMPLATE_DIR}/customer_template.pptx"
output_path = f"{OUTPUT_DIR}/proposal_output_test.pptx"
slide_keys = ["cover_page", "market_analysis_market_overview", "competitive_benchmarking"]

generate_ppt_from_template(slide_keys, template_path, output_path)
